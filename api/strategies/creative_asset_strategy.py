from abc import ABC, abstractmethod
from ctypes import Array
from pptx.util import Inches, Pt
from os import path, chdir, getcwd, mkdir
import requests
from flask_restful import abort
import numpy as np

# todo: add type hints everywhere
class CreativeAssetStrategy(ABC):
    """Base class for creative asset strategies."""

    def __init__(self, detail_slide: object, creative_asset_data: dict) -> None:
        self.detail_slide = detail_slide
        self.creative_asset_data = creative_asset_data

    @abstractmethod
    def create_creative_asset_title_textbox(self) -> None:
        pass

    @staticmethod
    def download_creative_asset(creative_asset_details: dict) -> str:
        """Download creative asset from Air Table link and cache locally.

        Downloads creative asset with `requests.get` and caches it in the
        `api/videos` or `api/images` folders depending on the asset's fiile extension.
        """
        link_to_asset = creative_asset_details["url"]
        filename = link_to_asset.split("/")[-1]
        file_extension = filename.split(".")[1]
        # todo: make types a property of the base class
        accepted_image_types = ["jpg", "png", "gif", "raw", "svg", "heic"]
        accepted_video_types = ["mp4", "mov", "m4v", "mpg", "mpeg", "wmv"]
        current_dir = path.abspath(getcwd())

        if path.isdir("api/videos") is False:
            mkdir("api/videos")
        if path.isdir("api/images") is False:
            mkdir("api/images")

        if file_extension in accepted_image_types:
            chdir("api/images")
            cached_file_path: str = path.abspath(filename)
        elif file_extension in accepted_video_types:
            chdir("api/videos")
            cached_file_path: str = path.abspath(filename)

        data = requests.get(link_to_asset)

        with open(cached_file_path, "wb") as f:
            f.write(data.content)
            chdir(current_dir)

        return cached_file_path

    def insert_creative(self) -> None:
        pass

    def build_asset_layout(self) -> None:
        pass


class SingleAsset(CreativeAssetStrategy):
    """Strategy for laying out slide with single creative asset."""

    # inherited init includes: detail_slide, creative_asset_data: dict

    # Done!
    def create_creative_asset_title_textbox(self) -> None:
        """Create creative asset textbox as default."""
        # Text box positioning and size definitions
        # Insert creative asset as normal
        left = Inches(0.33)
        top = Inches(4.58)
        width = Inches(4.4)
        height = Inches(0.4)

        slide = self.detail_slide

        # Textbox setup
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text with run
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Creative Asset:"
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)

    def insert_creative(self, path_to_asset: str):
        # Add creative asset file
        filename = path.basename(path_to_asset)
        file_extension = filename.split(".")[1]

        # throw error if pdf extension
        if file_extension == "pdf":
            return abort(404, "pdf is not an accepted Asset type. Please convert PDF to JPG")

        accepted_image_types = [
            "jpg",
            "png",
            "gif",
            "raw",
            "svg",
            "heic",
        ]
        accepted_video_types = ["mp4", "mov", "m4v", "mpg", "mpeg", "wmv"]

        # todo: throw error if asset doesn't match accepted types

        slide = self.detail_slide

        if file_extension in accepted_image_types:
            # Creative asset positioning and size definitions
            left = Inches(0.43)
            top = Inches(5.13)
            width = Inches(4.1)
            slide.shapes.add_picture(path_to_asset, left, top, width)
        elif file_extension in accepted_video_types:
            # todo: add video thumbnail
            # Creative asset positioning and size definitions
            left = Inches(0.43)
            top = Inches(5.13)
            width = Inches(2.87)
            height = Inches(2.17)
            slide.shapes.add_movie(path_to_asset, left, top, width, height)

    def build_asset_layout(self) -> None:
        # create title textbox
        self.create_creative_asset_title_textbox()
        # download creative asset - pass in only single post
        creative_asset_details: dict = self.creative_asset_data[0]
        cached_file_path = self.download_creative_asset(creative_asset_details)
        # insert creative asset - pass in path to file
        self.insert_creative(cached_file_path)
        # ? need to return the modified detail slide?
        # return self.detail_slide


class AssetTrio(CreativeAssetStrategy):
    """Strategy for laying out slide with multiple creative assets."""

    # todo: need to fix params across strategies - should be list
    def __init__(self, detail_slide: object, creative_asset_data: list) -> None:
        self.detail_slide = detail_slide
        self.creative_asset_data = creative_asset_data

    def create_creative_asset_title_textbox(self) -> None:
        """Create creative asset textbox as default."""
        # Text box positioning and size definitions
        left = Inches(0.33)
        top = Inches(4.49)
        width = Inches(4.4)
        height = Inches(0.4)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text with run
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Creative Asset:"
        bold_run_font = bold_run.font
        bold_run_font.bold = True

    @property
    def accepted_image_types(self) -> list:
        return ["jpg", "png", "gif", "raw", "svg", "heic"]

    @property
    def accepted_video_types(self) -> list:
        return ["mp4", "mov", "m4v", "mpg", "mpeg", "wmv"]

    # @property
    # def accepted_extensions(self) -> list:
    #     return self.accepted_image_types + self.accepted_video_types

    #! using 2d matrices to reduce time complexity
    @property
    def size_and_position_matrix(self) -> Array:
        # need left, top, width, and height
        # need three rows
        #! may only want to specify width here to allow for vertical aspect ratios automatically
        image_1 = [0.42, 5.09, 2.74, 1.83]
        image_2 = [3.48, 5.09, 2.74, 1.83]
        image_3 = [6.54, 5.09, 2.74, 1.83]

        return np.array([image_1, image_2, image_3])

    # @staticmethod
    def extract_filename(self, asset_details: dict) -> str:
        return asset_details["filename"]

    # @staticmethod
    def extract_extension(self, filename: str) -> str:
        return filename.split(".")[1]

    # @staticmethod
    def download_creative_asset(self, asset_url: str) -> str:
        """Download creative asset from Air Table link and cache locally.

        Downloads creative asset with `requests.get` and caches it in the
        `api/videos` or `api/images` folders depending on the asset's fiile extension.
        """

        filename = asset_url.split("/")[-1]
        file_extension = filename.split(".")[1]
        # todo: make types a property of the base class
        accepted_image_types = ["jpg", "png", "gif", "raw", "svg", "heic"]
        accepted_video_types = ["mp4", "mov", "m4v", "mpg", "mpeg", "wmv"]
        current_dir = path.abspath(getcwd())

        if path.isdir("api/videos") is False:
            mkdir("api/videos")
        if path.isdir("api/images") is False:
            mkdir("api/images")

        if file_extension in accepted_image_types:
            chdir("api/images")
            cached_file_path: str = path.abspath(filename)
        elif file_extension in accepted_video_types:
            chdir("api/videos")
            cached_file_path: str = path.abspath(filename)

        data = requests.get(asset_url)

        with open(cached_file_path, "wb") as f:
            f.write(data.content)
            chdir(current_dir)

        return cached_file_path

    def add_creative_asset_image(self, path_to_asset: str, index: int) -> None:

        slide = self.detail_slide

        left = Inches(self.size_and_position_matrix[index][0])
        top = Inches(self.size_and_position_matrix[index][1])
        width = Inches(self.size_and_position_matrix[index][2])
        slide.shapes.add_picture(path_to_asset, left, top, width)

    def add_creative_asset_video(self, path_to_asset: str, index: int) -> None:

        slide = self.detail_slide

        # todo: use pixel counts from the asset dictionary to assign these
        left = Inches(self.size_and_position_matrix[index][0])
        top = Inches(self.size_and_position_matrix[index][1])
        width = Inches(self.size_and_position_matrix[index][2])
        height = Inches(self.size_and_position_matrix[index][3])
        slide.shapes.add_movie(path_to_asset, left, top, width, height)

    def build_asset_layout(self) -> None:

        self.create_creative_asset_title_textbox()
        for index, asset in enumerate(self.creative_asset_data):
            asset_url = asset["url"]
            filename = self.extract_filename(asset)
            extension = self.extract_extension(filename)
            cache_path = self.download_creative_asset(asset_url)
            if extension in self.accepted_image_types:
                self.add_creative_asset_image(cache_path, index)
            elif extension in self.accepted_video_types:
                self.add_creative_asset_video(cache_path, index)
