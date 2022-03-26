from abc import ABC, abstractmethod
import re
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from api.strategies.creative_asset_strategy import SingleAsset, MultipleAssets


class SlideLayout(ABC):
    def __init__(self, presentation: object, post_data: dict) -> None:
        self.presentation = presentation
        self.post_data = post_data
        self.platform: str = post_data["platform"]
        self.account: str = post_data["account"]
        self.country: str = post_data["country"]
        self.paid_v_organic = post_data["paid_v_organic"]
        self.audience: str = post_data["audience"]
        self.publish_date: str = self.format_publish_date(post_data["publish_date"])
        self.social_copy: str = post_data["copy"]
        self.link: str = post_data["link"]

    @classmethod
    def format_publish_date(self, publish_date: str) -> str:
        """Remove time from ISO formatted string."""

        # datetime implementation
        publish_dt: datetime = datetime.strptime(publish_date, r"%Y-%m-%dT%H:%M:%S.%fZ").date()
        dt_str: str = publish_dt.strftime(r"%m/%d/%Y")

        return dt_str

    def create_blank_slide(self) -> None:
        """Create blank power point slide.

        Selects a slide layout from avoialable power point options and
        adds a new slide with that layout to the passed instance of
        :param:presenation.

        The `Presentation` instance is inherited from base class `SlideLayout`.

        :return: returns slide object
        :rtype: `object`
        """

        blank_slide_layout = self.presentation.slide_layouts[6]
        detail_slide = self.presentation.slides.add_slide(blank_slide_layout)
        self.detail_slide = detail_slide

    def create_platform_textbox(self) -> None:
        """Create a textbox for a social platform."""
        # Text box positioning and size definitions
        left = Inches(0.33)
        top = Inches(0.21)
        width = Inches(1.68)
        height = Inches(0.34)
        # Textbox setup and adding text
        slide = self.detail_slide
        platform_tbox = slide.shapes.add_textbox(left, top, width, height)
        platform_tf = platform_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = platform_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Platform: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        bold_run_font.color = RGBColor(255, 255, 255)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.platform
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)
        normal_run_font.color = RGBColor(255, 255, 255)

    def create_country_textbox(self) -> None:
        """Create textbox for country."""
        # Text box positioning and size definitions
        left = Inches(3.08)
        top = Inches(0.21)
        width = Inches(1.17)
        height = Inches(0.34)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Country: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.country
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    def create_paid_v_organic_textbox(self) -> None:
        """Create textbox for designation of paid or organic."""
        # Text box positioning and size definitions
        left = Inches(5.54)
        top = Inches(0.21)
        width = Inches(1.51)
        height = Inches(0.34)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Type: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.paid_v_organic
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    def create_account_textbox(self) -> None:
        """Create a textbox for a social account."""
        # Text box positioning and size definitions
        left = Inches(0.33)
        top = Inches(0.71)
        width = Inches(1.81)
        height = Inches(0.34)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Account: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.account
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    def create_audience_textbox(self) -> None:
        """Create textbox for audience."""
        # Text box positioning and size definitions
        left = Inches(3.08)
        top = Inches(0.71)
        width = Inches(1.51)
        height = Inches(0.34)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Audience: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.audience
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    def create_publish_date_textbox(self) -> None:
        """Create textbox for publish date."""
        # Text box positioning and size definitions
        left = Inches(5.54)
        top = Inches(0.71)
        width = Inches(4.1)
        height = Inches(0.34)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Tentative Post Date: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        normal_run = p.add_run()
        normal_run.text = self.publish_date
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    @abstractmethod
    def create_social_copy_textbox(self) -> None:
        pass

    @abstractmethod
    def create_link_textbox(self) -> None:
        """Create textbox for link."""
        pass

    def add_creative_assets(self) -> None:

        # use creative asset strategy here
        if "creative" in self.post_data.keys():
            creative_asset_data: list = self.post_data["creative"]
            creative_asset_count = len(creative_asset_data)
            if creative_asset_count == 1:
                creative_strategy = SingleAsset(self.detail_slide, creative_asset_data)
                creative_strategy.build_asset_layout()

            elif (creative_asset_count > 1) and (creative_asset_count <= 12):
                # for asset in creative_asset_data:
                #     asset_counter = 1
                # InsertMultipleAssets(self.detail_slide, asset)
                pass
            else:
                # handle outlier cases here - maybe with error when number of creative assets exceeds 12
                pass
        else:
            # logic for posts with no creative asset
            pass

    @abstractmethod
    def build_layout(self) -> None:
        pass


class Organic(SlideLayout):
    """Strategy for creating organic post slide layout.

    Organic option for slide layouts. This format omits values that exist in other type of posts such as paid promotions. For example, this layout does not include a `paid_headline`.
    """

    def create_detail_background(self) -> None:
        slide = self.detail_slide
        left = top = Inches(0)
        width = Inches(10)
        height = Inches(1.33)
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rectangle_fill = rectangle.fill
        rectangle_fill.solid()
        rectangle_fill.fore_color.rgb = RGBColor(55, 96, 146)
        rectangle_outline = rectangle.line
        rectangle_outline.color.rgb = RGBColor(55, 96, 146)
        rectangle.shadow.inherit = False

    def create_social_copy_textbox(self) -> None:
        """Create textbox for social copy."""
        # Text box positioning and size definitions
        left = Inches(0.25)
        top = Inches(1.65)
        width = Inches(4.5)
        height = Inches(4.64)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        account_tf.word_wrap = True
        # Adding bold text and normal text in separate paragraphs
        p1 = account_tf.paragraphs[0]
        # Bold paragraph
        run = p1.add_run()
        run.text = "Copy:\n"
        p1_font = run.font
        p1_font.bold = True
        p1_font.size = Pt(14)
        # Normal paragraph
        p2 = account_tf.add_paragraph()
        p2.text = self.social_copy
        p2.font.size = Pt(14)

    def create_link_textbox(self) -> None:
        """Create textbox for link."""
        # Text box positioning and size definitions
        left = Inches(5.25)
        top = Inches(1.65)
        width = Inches(4.39)
        height = Inches(0.4)
        # Textbox setup
        slide = self.detail_slide
        account_tbox = slide.shapes.add_textbox(left, top, width, height)
        account_tf = account_tbox.text_frame
        # Adding bold text and normal text in single paragraph with runs
        p = account_tf.paragraphs[0]
        # Bold run
        bold_run = p.add_run()
        bold_run.text = "Link: "
        bold_run_font = bold_run.font
        bold_run_font.bold = True
        bold_run_font.size = Pt(14)
        # Normal run
        # todo: if there is no link, don't add a hyperlink address
        normal_run = p.add_run()
        normal_run.text = self.link
        hyperlink = normal_run.hyperlink
        hyperlink.address = self.link
        normal_run_font = normal_run.font
        normal_run_font.size = Pt(14)

    # todo: move build layout logic to base class
    def build_layout(self):
        # create blank slide
        self.create_blank_slide()
        # create detail background rectangle
        self.create_detail_background()
        # create all textboxes
        self.create_platform_textbox()
        self.create_country_textbox()
        self.create_account_textbox()
        self.create_audience_textbox()
        self.create_paid_v_organic_textbox()
        self.create_publish_date_textbox()
        self.create_social_copy_textbox()
        self.create_link_textbox()
        # todo: add paid/organic type and tentative publish date textboxes here
        # add creative asset
        self.add_creative_assets()


# class Paid(SlideLayout):
#     """Strategy for creating paid post slide layout."""

#     def __init__(self, presentation: object, post_data: dict) -> None:
#         super().__init__(presentation, post_data)
#         # add paid healdine property
#         self.paid_headline = post_data["paid_headline"]

#     # new paid headline copy section - social copy moved down to accomodate
#     def create_paid_headline_textbox(self) -> None:
#         """Create textbox for social copy."""
#         # Text box positioning and size definitions
#         left = Inches(0.25)
#         top = Inches(1.65)
#         width = Inches(4.5)
#         height = Inches(4.64)
#         # Textbox setup
#         slide = self.detail_slide
#         account_tbox = slide.shapes.add_textbox(left, top, width, height)
#         account_tf = account_tbox.text_frame
#         account_tf.word_wrap = True
#         # Adding bold text and normal text in separate paragraphs
#         p1 = account_tf.paragraphs[0]
#         # Bold paragraph
#         run = p1.add_run()
#         run.text = "Copy:\n"
#         font = run.font
#         font.bold = True
#         # Normal paragraph
#         p2 = account_tf.add_paragraph()
#         p2.text = self.social_copy

#     def create_social_copy_textbox(self) -> None:
#         """Create textbox for social copy."""
#         # todo: update this based on paid specifications
#         # Text box positioning and size definitions
#         left = Inches(0.25)
#         top = Inches(1.65)
#         width = Inches(4.5)
#         height = Inches(4.64)
#         # Textbox setup
#         slide = self.detail_slide
#         account_tbox = slide.shapes.add_textbox(left, top, width, height)
#         account_tf = account_tbox.text_frame
#         account_tf.word_wrap = True
#         # Adding bold text and normal text in separate paragraphs
#         p1 = account_tf.paragraphs[0]
#         # Bold paragraph
#         run = p1.add_run()
#         run.text = "Copy:\n"
#         p1_font = run.font
#         p1_font.bold = True
#         p1_font.size = Pt(14)
#         # Normal paragraph
#         p2 = account_tf.add_paragraph()
#         p2.text = self.social_copy
#         p2.font.size = Pt(14)

#     # todo: implement build layout method
#     def build_layout(self) -> None:
#         # create blank slide per normal
#         self.create_blank_slide()
#         # create all textboxes
#         self.create_platform_textbox()
#         self.create_country_textbox()
#         self.create_account_textbox()
#         self.create_audience_textbox()
#         self.create_paid_headline_textbox()
#         self.create_social_copy_textbox()
#         self.create_link_textbox()
#         # todo: add paid/organic type and tentative publish date textboxes here
#         # add creative asset
#         self.add_creative_assets()
