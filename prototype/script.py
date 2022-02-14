# Script that contains rudamentary logic for ppt generator

# Import ppt package
from pptx import Presentation
from pptx.util import Inches
from os import path
import requests

# Todo: Add Paid/Organic and Paid Headline elements
# Todo: Enable multiple creative assets to be added for each post if required


# Slide setup
def create_blank_slide(prs):
    # Select slide layout and create a new slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide


def create_platform_textbox(slide, platform: str):
    # Text box positioning and size definitions
    left = Inches(0.25)
    top = Inches(0.25)
    width = Inches(4.5)
    height = Inches(0.4)
    # Textbox setup and adding text
    platform_tbox = slide.shapes.add_textbox(left, top, width, height)
    platform_tf = platform_tbox.text_frame
    # Adding bold text and normal text in single paragraph with runs
    p = platform_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Platform: "
    bold_run_font = bold_run.font
    bold_run_font.bold = True
    # Normal run
    normal_run = p.add_run()
    normal_run.text = platform


def create_account_textbox(slide, account: str):
    # Text box positioning and size definitions
    left = Inches(0.25)
    top = Inches(0.75)
    width = Inches(4.5)
    height = Inches(0.4)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    # Adding bold text and normal text in single paragraph with runs
    p = account_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Account: "
    bold_run_font = bold_run.font
    bold_run_font.bold = True
    # Normal run
    normal_run = p.add_run()
    normal_run.text = account


def create_country_textbox(slide, country: str):
    # Text box positioning and size definitions
    left = Inches(5.25)
    top = Inches(0.25)
    width = Inches(4.39)
    height = Inches(0.4)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    # Adding bold text and normal text in single paragraph with runs
    p = account_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Country: "
    bold_run_font = bold_run.font
    bold_run_font.bold = True
    # Normal run
    normal_run = p.add_run()
    normal_run.text = country


def create_audience_textbox(slide, audience: str):
    # Text box positioning and size definitions
    left = Inches(5.25)
    top = Inches(0.75)
    width = Inches(4.39)
    height = Inches(0.4)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    # Adding bold text and normal text in single paragraph with runs
    p = account_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Country: "
    bold_run_font = bold_run.font
    bold_run_font.bold = True
    # Normal run
    normal_run = p.add_run()
    normal_run.text = audience


def create_social_copy_textbox(slide, social_copy: str):
    # Text box positioning and size definitions
    left = Inches(0.25)
    top = Inches(1.65)
    width = Inches(4.5)
    height = Inches(4.64)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    account_tf.word_wrap = True
    # Adding bold text and normal text in separate paragraphs
    p1 = account_tf.paragraphs[0]
    # Bold paragraph
    run = p1.add_run()
    run.text = "Copy:\n"
    font = run.font
    font.bold = True
    # Normal paragraph
    p2 = account_tf.add_paragraph()
    p2.text = social_copy


def create_link_textbox(slide, link: str):
    # Text box positioning and size definitions
    left = Inches(5.25)
    top = Inches(1.65)
    width = Inches(4.39)
    height = Inches(0.4)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    # Adding bold text and normal text in single paragraph with runs
    p = account_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Link: "
    bold_run_font = bold_run.font
    bold_run_font.bold = True
    # Normal run
    normal_run = p.add_run()
    normal_run.text = link
    hyperlink = normal_run.hyperlink
    hyperlink.address = link


def create_creative_asset_textbox(slide):
    # Text box positioning and size definitions
    left = Inches(5.25)
    top = Inches(2.55)
    width = Inches(4.4)
    height = Inches(0.4)
    # Textbox setup
    account_tbox = slide.shapes.add_textbox(left, top, width, height)
    account_tf = account_tbox.text_frame
    # Adding bold text with run
    p = account_tf.paragraphs[0]
    # Bold run
    bold_run = p.add_run()
    bold_run.text = "Creative Asset:"
    bold_run_font = bold_run.font
    bold_run_font.bold = True


def download_creative_asset(link_to_asset: str) -> str:
    filename = link_to_asset.split("/")[-1]
    cached_file_path = filename
    data = requests.get(link_to_asset)
    with open(cached_file_path, "wb") as f:
        f.write(data.content)
    return cached_file_path


# ? --- Use context manager to run this function? ---
def insert_creative_asset(slide, path_to_asset: str):
    # Add creative asset file
    filename = path.basename(path_to_asset)
    file_extension = filename.split(".")[1]
    accepted_image_types = ["jpg", "png", "gif", "raw", "svg", "heic"]
    accepted_video_types = ["mp4", "mov", "m4v", "mpg", "mpeg", "wmv"]
    if file_extension in accepted_image_types:
        # Creative asset positioning and size definitions
        left = Inches(5.25)
        top = Inches(3.1)
        width = Inches(4.1)
        slide.shapes.add_picture(path_to_asset, left, top, width)
    elif file_extension in accepted_video_types:
        # Creative asset positioning and size definitions
        left = Inches(5.25)
        top = Inches(3.1)
        width = Inches(4.1)
        height = Inches(3.1)
        slide.shapes.add_movie(path_to_asset, left, top, width, height)


# Todo: Create crud operations for handling cached image assets


# Main function that runs script


def create_presentation():
    # Create a presentation
    prs = Presentation()
    return prs


def create_slide(prs, post_data):
    # Create blank slide
    slide = create_blank_slide(prs)
    # Add platform element to slide
    create_platform_textbox(slide, post_data["platform"])
    # Add account element to slide
    create_account_textbox(slide, post_data["account"])
    # Add country element
    create_country_textbox(slide, post_data["country"])
    # Add audience element
    create_audience_textbox(slide, post_data["audience"])
    # Add social copy element
    create_social_copy_textbox(slide, post_data["copy"])
    # Add link element
    create_link_textbox(slide, post_data["link"])
    # Add creative asset element
    create_creative_asset_textbox(slide)
    # Download creative asset
    link_to_asset = post_data["creative_asset"][0]["url"]
    cached_asset_path = download_creative_asset(link_to_asset)
    # Add creative asset
    path_to_asset = path.abspath(cached_asset_path)
    insert_creative_asset(slide, path_to_asset)


def save_ppt(prs, ppt_filename):
    # Save ppt file
    prs.save(f"{ppt_filename}.pptx")


if __name__ == "__main__":
    print("creating ppt")
    # Dummy Data
    dummy_data = {
        "group_name": "Brian Levitt: Bloomberg The Close Media Appearance",
        "posts": [
            {
                "platform": "LinkedIn",
                "account": "Brian Levitt",
                "country": "US",
                "audience": "Retail",
                "copy": "Are we approaching the end of a cycle or just navigating through a difficult period? Skip ahead to Global Market Strategist Brian Levitt’s interview with Romaine Bostick, Caroline Hyde, and Taylor Riggs at the 1:19:42 mark to hear his explanation for why he expects inflationary pressures to moderate over the next year.",
                "link": "https://inves.co/3HXoVO8",
                "creative_asset": [
                    {
                        "id": "att89hPO5Bb7nTSyT",
                        "width": 375,
                        "height": 196,
                        "url": "https://dl.airtable.com/.attachments/6480ee015e00f8e38dd1473aff3d1c34/d49b6910/BLYearInReviewSlide2.jpg",
                        "filename": "BL Year In Review Slide 2.jpg",
                        "size": 12942,
                        "type": "image/jpeg",
                        "thumbnails": {
                            "small": {"url": "https://dl.airtable.com/.attachmentThumbnails/f78da1a12ab9c9be87b9d279b091fa7f/a25ae696", "width": 69, "height": 36},
                            "large": {"url": "https://dl.airtable.com/.attachmentThumbnails/f9ddf156ea2c891522b9c44fc760bb8e/a37f3063", "width": 375, "height": 196},
                            "full": {"url": "https://dl.airtable.com/.attachmentThumbnails/6df0d1773d8d33b165c104e000de27bc/4b08459b", "width": 3000, "height": 3000},
                        },
                    },
                ],
            },
            {
                "platform": "LinkedIn",
                "account": "Brian Levitt",
                "country": "US",
                "audience": "Retail",
                "copy": "Are we approaching the end of a cycle or just navigating through a difficult period? Skip ahead to Global Market Strategist Brian Levitt’s interview with Romaine Bostick, Caroline Hyde, and Taylor Riggs at the 1:19:42 mark to hear his explanation for why he expects inflationary pressures to moderate over the next year.",
                "link": "https://inves.co/3HXoVO8",
                "creative_asset": [
                    {
                        "id": "attkntEYA67dwCIiK",
                        "width": 375,
                        "height": 196,
                        "url": "https://dl.airtable.com/.attachments/51fd5ccb3bd00df2815c90b6c1fd2692/402f9dfb/BLYearInReviewSlide3.jpg",
                        "filename": "BL Year In Review Slide 3.jpg",
                        "size": 16674,
                        "type": "image/jpeg",
                        "thumbnails": {
                            "small": {"url": "https://dl.airtable.com/.attachmentThumbnails/53863b92c733d492ba932f080333eae7/c7c48295", "width": 69, "height": 36},
                            "large": {"url": "https://dl.airtable.com/.attachmentThumbnails/eacadd4a1be4b00ebb6a6a8de5595318/44ce87d9", "width": 375, "height": 196},
                            "full": {"url": "https://dl.airtable.com/.attachmentThumbnails/42487d29d653e354e01a767695345ea6/4952131f", "width": 3000, "height": 3000},
                        },
                    }
                ],
            },
        ],
    }
    prs = create_presentation()
    for post in dummy_data["posts"]:
        create_slide(prs, post)
    save_ppt(prs, dummy_data["group_name"])
    print("ppt created")
